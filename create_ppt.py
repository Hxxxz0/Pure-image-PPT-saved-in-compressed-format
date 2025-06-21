import os
import glob
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import re
import json

def get_user_config():
    """获取用户配置，优先从压缩脚本的配置中读取"""
    print("=" * 60)
    print("PPT创建工具 - 配置向导")
    print("=" * 60)
    
    # 默认配置
    config = {
        'output_prefix': '图片集',
        'target_ppt_size_mb': 18,
        'input_pattern': '*.png'
    }
    
    # 获取输出文件名前缀
    print("\n请输入PPT文件名称:")
    print("   例如: '我的演示文稿' 或 '项目汇报' 或 '会议展示'")
    
    output_name = input("   PPT文件名称: ").strip()
    if output_name:
        config['output_prefix'] = output_name
    else:
        print(f"   使用默认名称: {config['output_prefix']}")
    
    return config

def create_ppt_from_compressed_images(output_prefix="图片集"):
    """使用压缩后的JPEG图片创建PPT，确保文件小于20MB"""
    
    # 首先检查压缩文件夹是否存在
    compressed_dir = "compressed_for_ppt"
    if not os.path.exists(compressed_dir):
        print("未找到压缩文件夹，请先运行 'python compress_images.py' 压缩图片")
        return None
    
    # 获取所有压缩后的JPEG文件并按数字顺序排序
    jpg_files = glob.glob(os.path.join(compressed_dir, "*_optimized.jpg"))
    
    # 提取文件名中的数字并排序
    def extract_number(filename):
        # 提取文件名中的所有数字，使用最后一个数字排序
        numbers = re.findall(r'\d+', os.path.basename(filename))
        return int(numbers[-1]) if numbers else 0
    
    # 按数字顺序排序
    jpg_files.sort(key=extract_number)
    
    if not jpg_files:
        print("未找到压缩后的图片文件")
        print("请确保已运行 'python compress_images.py' 压缩图片")
        return None
    
    print(f"找到 {len(jpg_files)} 个压缩后的图片文件")
    
    # 获取第一张图片的尺寸来确定幻灯片比例
    first_image = jpg_files[0]
    with Image.open(first_image) as img:
        img_width, img_height = img.size
        aspect_ratio = img_width / img_height
    
    print(f"检测到图片比例: {img_width}x{img_height} (比例: {aspect_ratio:.3f})")
    
    # 创建新的PowerPoint演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸以匹配图片比例
    # 使用标准的16:9比例，适合大多数显示器
    if aspect_ratio > 1.5:  # 接近16:9
        slide_width = Inches(10)
        slide_height = Inches(5.625)  # 16:9比例
    else:  # 更接近4:3或正方形
        slide_width = Inches(10)
        slide_height = Inches(7.5)  # 4:3比例
    
    # 设置幻灯片尺寸
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    
    print(f"设置幻灯片尺寸: {slide_width.inches:.2f}x{slide_height.inches:.2f}英寸")
    
    # 为每个图片创建一张幻灯片
    for i, image_path in enumerate(jpg_files):
        print(f"正在处理: {os.path.basename(image_path)}")
        
        # 添加空白幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        try:
            # 图片从左上角开始，充满整个幻灯片
            left = Inches(0)
            top = Inches(0)
            
            # 添加图片到幻灯片，完全填满幻灯片区域
            slide.shapes.add_picture(image_path, left, top, slide_width, slide_height)
            
        except Exception as e:
            print(f"处理图片 {image_path} 时出错: {e}")
            continue
    
    # 保存PPT文件
    output_filename = f"{output_prefix}_压缩版.pptx"
    prs.save(output_filename)
    
    # 检查文件大小
    file_size = os.path.getsize(output_filename)
    file_size_mb = file_size / 1024 / 1024
    
    print(f"PPT已保存为: {output_filename}")
    print(f"文件大小: {file_size_mb:.2f} MB")
    print(f"总共创建了 {len(prs.slides)} 张幻灯片")
    
    if file_size_mb > 20:
        print("⚠️  警告：文件大小超过20MB，可能需要进一步压缩")
    else:
        print("✅ 文件大小符合要求（< 20MB）")
    
    return output_filename

def create_ppt_from_images(input_pattern="*.png", output_prefix="图片集"):
    """原版本：将PNG图片作为每张幻灯片的背景，保持图片原始比例"""
    
    # 获取所有图片文件并按数字顺序排序
    image_files = glob.glob(input_pattern)
    
    # 提取文件名中的数字并排序
    def extract_number(filename):
        numbers = re.findall(r'\d+', os.path.basename(filename))
        return int(numbers[-1]) if numbers else 0
    
    # 按数字顺序排序
    image_files.sort(key=extract_number)
    
    if not image_files:
        print(f"未找到匹配模式 '{input_pattern}' 的图片文件")
        print("请检查文件名模式是否正确")
        return None
    
    print(f"找到 {len(image_files)} 个图片文件")
    
    # 获取第一张图片的尺寸来确定幻灯片比例
    first_image = image_files[0]
    with Image.open(first_image) as img:
        img_width, img_height = img.size
        aspect_ratio = img_width / img_height
    
    print(f"检测到图片比例: {img_width}x{img_height} (比例: {aspect_ratio:.3f})")
    
    # 创建新的PowerPoint演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸以匹配图片比例
    # 以10英寸为基准宽度，高度根据比例计算
    slide_width = Inches(10)
    slide_height = Inches(10 / aspect_ratio)
    
    # 设置幻灯片尺寸
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    
    print(f"设置幻灯片尺寸: {slide_width.inches:.2f}x{slide_height.inches:.2f}英寸")
    
    # 为每个图片创建一张幻灯片
    for i, image_path in enumerate(image_files):
        print(f"正在处理: {os.path.basename(image_path)}")
        
        # 添加空白幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        try:
            # 图片从左上角开始，充满整个幻灯片
            left = Inches(0)
            top = Inches(0)
            
            # 添加图片到幻灯片，完全填满幻灯片区域
            slide.shapes.add_picture(image_path, left, top, slide_width, slide_height)
            
        except Exception as e:
            print(f"处理图片 {image_path} 时出错: {e}")
            continue
    
    # 保存PPT文件
    output_filename = f"{output_prefix}_原始版.pptx"
    prs.save(output_filename)
    print(f"PPT已保存为: {output_filename}")
    print(f"总共创建了 {len(prs.slides)} 张幻灯片")
    
    return output_filename

if __name__ == "__main__":
    print("欢迎使用PPT创建工具！")
    print("本工具可以将图片转换为PPT演示文稿。")
    print()
    
    # 获取用户配置
    config = get_user_config()
    
    print("\n选择PPT创建模式:")
    print("1. 使用压缩图片（推荐，文件小于20MB）")
    print("2. 使用原始图片（文件可能很大）")
    
    choice = input("\n请输入选择 (1/2): ").strip()
    
    print("\n" + "=" * 60)
    print("开始创建PPT...")
    print("=" * 60)
    
    if choice == "1":
        result = create_ppt_from_compressed_images(config['output_prefix'])
    elif choice == "2":
        print("\n请输入原始图片文件名模式:")
        print("例如: '*.png' 或 'slide_*.jpg' 或 'presentation_*.png'")
        pattern = input("文件名模式: ").strip()
        if not pattern:
            pattern = "*.png"
        result = create_ppt_from_images(pattern, config['output_prefix'])
    else:
        print("使用默认选项：压缩版本")
        result = create_ppt_from_compressed_images(config['output_prefix'])
    
    if result:
        print("\n" + "=" * 60)
        print(f"🎉 PPT创建成功！文件已保存为: {result}")
        print("=" * 60)
    else:
        print("\n❌ PPT创建失败，请检查配置和文件") 